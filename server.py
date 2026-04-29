"""
server.py  —  FastAPI 後端
端點：
    GET  /templates                  → 範本清單 + PDF 預覽 URL
    GET  /previews/{filename}        → 範本 PDF 預覽（靜態）
    POST /generate-pptx-preview      → AI 內容 → pptx → PDF，回傳預覽 URL + job_id
    GET  /download/{job_id}          → 下載已產生的 pptx
"""

from fastapi import FastAPI, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, FileResponse
from pydantic import BaseModel
from typing import Optional, List, Dict
from pathlib import Path
from urllib.parse import quote, quote as url_quote
import subprocess, uuid, io, tempfile, shutil, urllib.request, urllib.parse, json as _json
import logging, datetime, math as _math, re as _re

from render_ppt import detect_layouts, render_slide
from pptx import Presentation
import fitz  # PyMuPDF — PDF → PNG，不依賴 pdftoppm

TEMPLATE_DIR = Path('./input')
PREVIEW_DIR  = Path('./cache/previews')   # 範本預覽 PDF
JOBS_DIR     = Path('./cache/jobs')       # 每次生成的暫存 pptx / pdf

# ── Debug Logger ─────────────────────────────────────────────────────────────
LOG_DIR  = Path('./cache/logs')
LOG_DIR.mkdir(parents=True, exist_ok=True)
LOG_FILE = LOG_DIR / 'debug.log'

logging.basicConfig(
    level    = logging.DEBUG,
    format   = '%(asctime)s [%(levelname)s] %(message)s',
    datefmt  = '%Y-%m-%d %H:%M:%S',
    handlers = [
        logging.FileHandler(LOG_FILE, encoding='utf-8'),
        logging.StreamHandler(),   # 同時印到 terminal
    ]
)
log = logging.getLogger('server')

# Load HF API key for image generation MCP tool
HF_API_KEY: str = ''
_hf_cfg = Path('./config/api_key.config')
if _hf_cfg.exists():
    for _ln in _hf_cfg.read_text(encoding='utf-8').splitlines():
        _ln = _ln.strip()
        if _ln.startswith('HF_API_KEY='):
            HF_API_KEY = _ln.split('=', 1)[1].strip()
            break

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["GET", "POST"],
    allow_headers=["*"],
)

# ── Schemas ──
class SlideData(BaseModel):
    type: str
    title: Optional[str] = ''
    subtitle: Optional[str] = None
    bullets: Optional[List[str]] = None
    body: Optional[str] = None
    quote: Optional[str] = None

class GenerateRequest(BaseModel):
    topic: str
    theme: Optional[str] = 'midnight'
    template: Optional[str] = None
    slides: List[SlideData]

# ── Layout cache ──
_layout_cache: Dict = {}

def get_layout_map(template_name: str) -> Dict:
    if template_name not in _layout_cache:
        path = TEMPLATE_DIR / template_name
        if not path.exists():
            raise FileNotFoundError(f"找不到範本：{template_name}")
        _layout_cache[template_name] = detect_layouts(str(path))
    return _layout_cache[template_name]


# ══════════════════════════════════════════════
#  soffice 轉 PDF 工具函數
# ══════════════════════════════════════════════

def pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Optional[Path]:
    """pptx → pdf"""
    out_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        ['soffice', '--headless', '--convert-to', 'pdf',
         '--outdir', str(out_dir), str(pptx_path)],
        capture_output=True, text=True, timeout=60
    )
    if result.returncode != 0:
        print(f"⚠️  soffice 失敗：{result.stderr}")
        return None
    pdf_path = out_dir / (pptx_path.stem + '.pdf')
    return pdf_path if pdf_path.exists() else None


def pdf_to_thumbnail(pdf_path: Path, out_path: Path, dpi: int = 150) -> bool:
    """PDF 第一頁 → PNG 縮圖（使用 PyMuPDF，不依賴 pdftoppm）"""
    log.debug(f'pdf_to_thumbnail: {pdf_path} → {out_path}  dpi={dpi}')
    try:
        doc  = fitz.open(str(pdf_path))
        page = doc[0]
        mat  = fitz.Matrix(dpi / 72, dpi / 72)   # 72 是 PDF 預設 DPI
        pix  = page.get_pixmap(matrix=mat, alpha=False)
        pix.save(str(out_path))
        doc.close()
        log.debug(f'pdf_to_thumbnail: 成功，檔案大小={out_path.stat().st_size} bytes')
        return True
    except Exception as e:
        log.error(f'pdf_to_thumbnail 失敗: {e}')
        return False


def pdf_to_slides(pdf_path: Path, out_dir: Path, dpi: int = 150) -> list[Path]:
    """PDF 所有頁 → PNG 列表（使用 PyMuPDF，取代 pdftoppm）"""
    log.debug(f'pdf_to_slides: {pdf_path} → {out_dir}  dpi={dpi}')
    out_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    try:
        doc = fitz.open(str(pdf_path))
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for i, page in enumerate(doc):
            pix      = page.get_pixmap(matrix=mat, alpha=False)
            out_path = out_dir / f'slide-{i + 1}.png'
            pix.save(str(out_path))
            paths.append(out_path)
            log.debug(f'  第 {i+1} 頁 → {out_path.name}')
        doc.close()
        log.info(f'pdf_to_slides: 共產生 {len(paths)} 張投影片圖片')
    except Exception as e:
        log.error(f'pdf_to_slides 失敗: {e}')
    return paths


# ══════════════════════════════════════════════
#  Startup：預先產生所有範本的 PDF 預覽
# ══════════════════════════════════════════════

@app.on_event("startup")
def startup():
    JOBS_DIR.mkdir(parents=True, exist_ok=True)
    if not TEMPLATE_DIR.exists():
        print("⚠️  input/ 資料夾不存在"); return

    for pptx in TEMPLATE_DIR.glob("*.pptx"):
        pdf   = PREVIEW_DIR / (pptx.stem + '.pdf')
        thumb = PREVIEW_DIR / (pptx.stem + '.png')
        stale = not pdf.exists() or pdf.stat().st_mtime < pptx.stat().st_mtime
        if stale:
            log.info(f'[startup] 產生預覽 PDF：{pptx.name}')
            pdf = pptx_to_pdf(pptx, PREVIEW_DIR)
        if pdf and (stale or not thumb.exists()):
            log.info(f'[startup] 產生縮圖：{pptx.stem}.png')
            pdf_to_thumbnail(pdf, thumb)
        try:
            get_layout_map(pptx.name)
        except Exception as e:
            log.warning(f'[startup] layout 偵測失敗 {pptx.name}：{e}')
    log.info('[startup] ✅ 啟動完成')


# ══════════════════════════════════════════════
#  GET /templates
# ══════════════════════════════════════════════

@app.get("/templates")
def list_templates():
    if not TEMPLATE_DIR.exists():
        return {"templates": []}
    result = []
    for pptx in sorted(TEMPLATE_DIR.glob("*.pptx")):
        thumb = PREVIEW_DIR / (pptx.stem + '.png')
        result.append({
            "name"     : pptx.name,
            "label"    : pptx.stem,
            "thumbnail": f"/thumbnails/{pptx.stem}.png" if thumb.exists() else None,
        })
    return {"templates": result}


# ══════════════════════════════════════════════
#  GET /previews/{filename}
# ══════════════════════════════════════════════

@app.get("/previews/{filename}")
def get_preview(filename: str):
    if not filename.endswith('.pdf'):
        raise HTTPException(400, "只提供 PDF 檔")
    path = PREVIEW_DIR / filename
    if not path.exists():
        raise HTTPException(404, "預覽不存在")
    return FileResponse(str(path), media_type="application/pdf")


@app.get("/thumbnails/{filename}")
def get_thumbnail(filename: str):
    if not filename.endswith('.png'):
        raise HTTPException(400, "只提供 PNG 檔")
    path = PREVIEW_DIR / filename
    if not path.exists():
        raise HTTPException(404, "縮圖不存在")
    return FileResponse(str(path), media_type="image/png")


# ══════════════════════════════════════════════
#  POST /generate-pptx-preview
#  → 產生 pptx + PDF，回傳 { preview_url, job_id }
# ══════════════════════════════════════════════

@app.post("/generate-pptx-preview")
def generate_pptx_preview(req: GenerateRequest):
    template_name = req.template or next(
        (p.name for p in sorted(TEMPLATE_DIR.glob("*.pptx"))), None
    )
    if not template_name:
        raise HTTPException(400, "input/ 內沒有任何範本")
    template_path = TEMPLATE_DIR / template_name
    if not template_path.exists():
        raise HTTPException(404, f"找不到範本：{template_name}")

    log.info(f'[generate] 開始 topic={req.topic!r} template={template_name} slides={len(req.slides)}')
    try:
        # 1. 產生 pptx
        layout_map = get_layout_map(template_name)
        prs = Presentation(str(template_path))
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        for sd in req.slides:
            layout_idx = layout_map.get(sd.type, 1)
            render_slide(prs, layout_idx, sd.model_dump())

        # 2. 存到 jobs/
        job_id   = str(uuid.uuid4())
        job_dir  = JOBS_DIR / job_id
        job_dir.mkdir(parents=True)
        pptx_path = job_dir / 'result.pptx'
        prs.save(str(pptx_path))

        # 3. 轉 PDF
        pdf_path = pptx_to_pdf(pptx_path, job_dir)
        if not pdf_path:
            raise RuntimeError("PDF 轉換失敗，請確認 soffice 是否正常")

        # 4. 產生每頁縮圖（PyMuPDF）
        log.info(f'[generate] job_id={job_id} 開始產生投影片縮圖')
        thumb_dir  = job_dir / 'thumbs'
        slide_imgs = pdf_to_slides(pdf_path, thumb_dir)
        log.info(f'[generate] 完成，共 {len(slide_imgs)} 張')

        return {
            "job_id"     : job_id,
            "preview_url": f"/job-preview/{job_id}",
        }

    except HTTPException:
        raise
    except Exception as e:
        log.exception(f'[generate] 失敗：{e}')
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════
#  GET /job-preview/{job_id}  — 預覽 PDF
# ══════════════════════════════════════════════

@app.get("/job-preview/{job_id}")
def job_preview(job_id: str):
    """回傳各頁縮圖的 URL 清單"""
    thumb_dir = JOBS_DIR / job_id / 'thumbs'
    if not thumb_dir.exists():
        raise HTTPException(404, "預覽不存在")
    slides = sorted(thumb_dir.glob("*.png"))
    urls = [f"/job-slide/{job_id}/{p.name}" for p in slides]
    return {"slides": urls}


@app.get("/job-slide/{job_id}/{filename}")
def job_slide(job_id: str, filename: str):
    path = JOBS_DIR / job_id / 'thumbs' / filename
    if not path.exists():
        raise HTTPException(404, "圖片不存在")
    return FileResponse(str(path), media_type="image/png")


# ══════════════════════════════════════════════
#  GET /download/{job_id}  — 下載 pptx
# ══════════════════════════════════════════════

# ══════════════════════════════════════════════
#  PATCH /job/{job_id}/slide/{idx}
#  更新單張投影片內容 → 重新產生 pptx + 該頁 PNG
# ══════════════════════════════════════════════

class SlideUpdate(BaseModel):
    slide: SlideData
    template: Optional[str] = None
    all_slides: List[SlideData]   # 完整 slides，重新產生整份 pptx

@app.patch("/job/{job_id}/slide/{slide_idx}")
def update_slide(job_id: str, slide_idx: int, body: SlideUpdate):
    job_dir   = JOBS_DIR / job_id
    pptx_path = job_dir / 'result.pptx'
    if not pptx_path.exists():
        raise HTTPException(404, "Job 不存在")

    template_name = body.template or next(
        (p.name for p in sorted(TEMPLATE_DIR.glob("*.pptx"))), None
    )
    template_path = TEMPLATE_DIR / template_name
    if not template_path.exists():
        raise HTTPException(404, f"找不到範本：{template_name}")

    try:
        # 用完整 slides 重建整份 pptx
        layout_map = get_layout_map(template_name)
        prs = Presentation(str(template_path))
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        for sd in body.all_slides:
            layout_idx = layout_map.get(sd.type, 1)
            render_slide(prs, layout_idx, sd.model_dump())
        prs.save(str(pptx_path))

        # 重新轉 PDF
        pdf_path = pptx_to_pdf(pptx_path, job_dir)
        if not pdf_path:
            raise RuntimeError("PDF 轉換失敗")

        # 重新產生所有縮圖（PyMuPDF）
        thumb_dir = job_dir / 'thumbs'
        for old_img in thumb_dir.glob("*.png"):
            old_img.unlink()
        log.info(f'[update_slide] job_id={job_id} slide_idx={slide_idx} 重新產生縮圖')
        pdf_to_slides(pdf_path, thumb_dir)

        # 回傳更新後的縮圖 URL 列表
        thumbs = sorted(thumb_dir.glob("*.png"))
        urls   = [f"/job-slide/{job_id}/{p.name}" for p in thumbs]
        log.info(f'[update_slide] 完成，共 {len(urls)} 張')
        return {"slides": urls, "updated_idx": slide_idx}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════
#  MCP-compatible Tool Server
#  GET  /mcp/tools       → 工具清單（MCP manifest）
#  POST /mcp/tools/call  → 執行工具（統一入口）
# ══════════════════════════════════════════════

MCP_TOOLS = [
    {
        "name": "get_datetime",
        "description": "取得目前的日期、時間與星期（台北時區）",
        "inputSchema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "calculate",
        "description": "計算數學運算式，支援 + - * / ^ sqrt abs sin cos tan log pi e 等",
        "inputSchema": {
            "type": "object",
            "properties": {
                "expression": {
                    "type": "string",
                    "description": "要計算的算式，例如 '2^10' 或 'sqrt(144) + pi'",
                },
            },
            "required": ["expression"],
        },
    },
    {
        "name": "search_web",
        "description": "透過 Wikipedia 搜尋主題資訊，回傳標題、摘要與連結。中文主題用 lang='zh'，英文主題用 lang='en'",
        "inputSchema": {
            "type": "object",
            "properties": {
                "query":       {"type": "string",  "description": "搜尋關鍵字"},
                "lang":        {"type": "string",  "description": "語言版本：zh（預設）或 en", "enum": ["zh", "en"]},
                "max_results": {"type": "number",  "description": "回傳筆數，預設 4，最多 8"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "generate_image",
        "description": "依照使用者的描述生成一張圖片（FLUX.1-schnell）。使用者說要「生成圖片」、「畫一張」、「generate an image」等時呼叫。",
        "inputSchema": {
            "type": "object",
            "properties": {
                "prompt": {
                    "type": "string",
                    "description": "詳細的圖片描述，英文效果較佳，例如 'a cute cat sitting on a wooden desk, photorealistic'",
                },
            },
            "required": ["prompt"],
        },
    },
]


class MCPCallRequest(BaseModel):
    name: str
    arguments: Dict = {}


@app.get("/mcp/tools")
def mcp_list_tools():
    """MCP 工具清單 — 供前端或外部 MCP client 查詢可用工具"""
    return {"tools": MCP_TOOLS}


@app.post("/mcp/tools/call")
def mcp_call_tool(req: MCPCallRequest):
    """MCP 工具呼叫 — 前端 tool use 的統一執行入口"""
    log.info(f'[MCP] call name={req.name!r} args={req.arguments}')

    if req.name == "get_datetime":
        tz_taipei = datetime.timezone(datetime.timedelta(hours=8))
        now       = datetime.datetime.now(tz=tz_taipei)
        weekdays  = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"]
        result    = {
            "datetime": now.strftime("%Y/%m/%d %H:%M:%S"),
            "weekday":  weekdays[now.weekday()],
            "iso":      now.isoformat(),
        }
        return {"content": [{"type": "text", "text": _json.dumps(result, ensure_ascii=False)}]}

    elif req.name == "calculate":
        expr = str(req.arguments.get("expression", ""))
        try:
            safe_expr = _re.sub(r'\^', r'**', expr)
            safe_ns   = {k: getattr(_math, k) for k in dir(_math) if not k.startswith('_')}
            safe_ns.update({"abs": abs, "round": round})
            value  = eval(safe_expr, {"__builtins__": {}}, safe_ns)  # noqa: S307
            result = {"expression": expr, "result": str(value)}
        except Exception as e:
            raise HTTPException(400, f"計算錯誤：{e}")
        return {"content": [{"type": "text", "text": _json.dumps(result, ensure_ascii=False)}]}

    elif req.name == "search_web":
        q           = req.arguments.get("query", "")
        lang        = req.arguments.get("lang", "zh")
        max_results = int(min(req.arguments.get("max_results", 4), 8))
        result      = search_web(q=q, max_results=max_results, lang=lang)
        return {"content": [{"type": "text", "text": _json.dumps(result, ensure_ascii=False)}]}

    elif req.name == "generate_image":
        if not HF_API_KEY:
            raise HTTPException(400, "HF_API_KEY 未設定，請在 config/api_key.config 加入 HF_API_KEY=hf_...")
        prompt = req.arguments.get("prompt", "").strip()
        if not prompt:
            raise HTTPException(400, "請提供 prompt")
        log.info(f'[MCP generate_image] prompt={prompt!r}')
        hf_url = "https://router.huggingface.co/hf-inference/models/black-forest-labs/FLUX.1-schnell"
        hf_req = urllib.request.Request(
            hf_url,
            data=_json.dumps({"inputs": prompt}).encode(),
            headers={"Authorization": f"Bearer {HF_API_KEY}", "Content-Type": "application/json"},
            method="POST",
        )
        try:
            with urllib.request.urlopen(hf_req, timeout=60) as r:
                img_bytes = r.read()
        except Exception as e:
            raise HTTPException(502, f"Hugging Face API 失敗：{e}")
        gen_dir  = Path('./cache/generated')
        gen_dir.mkdir(parents=True, exist_ok=True)
        img_name = f"{uuid.uuid4().hex[:8]}.png"
        (gen_dir / img_name).write_bytes(img_bytes)
        result = {"image_url": f"/generated/{img_name}", "prompt": prompt}
        return {"content": [{"type": "text", "text": _json.dumps(result, ensure_ascii=False)}]}

    else:
        raise HTTPException(404, f"未知工具：{req.name}")


# ══════════════════════════════════════════════
#  GET /search?q=...&max_results=4
#  後端代理 Wikipedia（避免瀏覽器 CORS 限制）
# ══════════════════════════════════════════════

@app.get("/search")
def search_web(q: str = Query(..., description="搜尋關鍵字"),
               max_results: int = Query(4, ge=1, le=8),
               lang: str = Query("zh", description="語言：zh（中文維基）或 en（英文維基）")):
    """
    搜尋後端：先查 Wikipedia 標題搜尋，再取各頁面的摘要。
    支援中文（lang=zh）與英文（lang=en）。
    """
    encoded  = urllib.parse.quote(q)
    wiki_api = f"https://{lang}.wikipedia.org/w/api.php"

    # 1. 搜尋關鍵字 → 取得相關頁面標題
    search_url = (
        f"{wiki_api}?action=query&list=search"
        f"&srsearch={encoded}&srlimit={max_results}"
        f"&format=json&utf8=1"
    )
    try:
        req = urllib.request.Request(search_url, headers={"User-Agent": "HW1-ChatGPT/1.0"})
        with urllib.request.urlopen(req, timeout=10) as r:
            search_data = _json.loads(r.read().decode())
    except Exception as e:
        raise HTTPException(502, f"Wikipedia 搜尋失敗：{e}")

    hits   = search_data.get("query", {}).get("search", [])
    titles = [h["title"] for h in hits[:max_results]]

    if not titles:
        return {"query": q, "lang": lang, "results": [], "note": "沒有找到相關結果"}

    # 2. 逐一取各頁面摘要
    results = []
    for title in titles:
        summary_url = (
            f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/"
            + urllib.parse.quote(title.replace(" ", "_"))
        )
        try:
            req2 = urllib.request.Request(summary_url, headers={"User-Agent": "HW1-ChatGPT/1.0"})
            with urllib.request.urlopen(req2, timeout=8) as r2:
                page = _json.loads(r2.read().decode())
            results.append({
                "title":   page.get("title", title),
                "snippet": page.get("extract", "")[:400],
                "url":     page.get("content_urls", {}).get("desktop", {}).get("page", ""),
                "source":  f"wikipedia-{lang}",
            })
        except Exception:
            continue  # 跳過取摘要失敗的頁面

    return {"query": q, "lang": lang, "results": results}


@app.get("/generated/{filename}")
def get_generated_image(filename: str):
    path = Path('./cache/generated') / filename
    if not path.exists():
        raise HTTPException(404, "圖片不存在")
    return FileResponse(str(path), media_type="image/png")


@app.get("/download/{job_id}")
def download_pptx(job_id: str):
    pptx = JOBS_DIR / job_id / 'result.pptx'
    if not pptx.exists():
        raise HTTPException(404, "檔案不存在（可能已過期）")
    return FileResponse(
        str(pptx),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="result.pptx",
    )