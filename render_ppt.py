"""
render_ppt.py
─────────────────────────────────────────────────────────────
流程：
  1. 讀取 JS 生成的 source.pptx，抽出每張投影片的結構化資料
  2. 讀取範本 wisptype.pptx，自動偵測每種 slide type 最適合的 layout
  3. 依序把內容填入範本，輸出 output/result.pptx

支援的 slide type：title / closing / bullets / content / quote
（對應 ppt.js 的 JSON 結構）
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import List, Dict, Optional
import re

# ── 範本路徑（server.py 也會 import 這個） ──
TEMPLATE_PATH = './input/wisptype.pptx'


# ══════════════════════════════════════════════
#  1. 字型 / 語言套用
# ══════════════════════════════════════════════

def detect_lang(text: str) -> str:
    """簡易偵測：含中文字元 → zh-TW，否則 en-US"""
    if re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', text):
        return 'zh-TW'
    return 'en-US'

def detect_font(lang: str) -> str:
    return 'Microsoft JhengHei' if lang == 'zh-TW' else 'Century Gothic'

def apply_style(run, text: str, font_name: str = None, font_size: int = None):
    """套用字型、語言、東亞字型標籤"""
    lang = detect_lang(text)
    font = font_name or detect_font(lang)

    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', lang)
    run.font.name = font
    if font_size:
        run.font.size = Pt(font_size)

    # a:latin
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        latin.set('typeface', font)

    # a:eaFont（東亞字型，確保中文正確渲染）
    ea = rPr.find(qn('a:eaFont')) or rPr.find(qn('a:ea'))
    if ea is None:
        ea = OxmlElement('a:eaFont')
        rPr.insert(0, ea)
    ea.set('typeface', font)


# ══════════════════════════════════════════════
#  2. 從 JS 生成的 pptx 抽取結構化資料
# ══════════════════════════════════════════════

def extract_slides(src_path: str) -> List:
    """
    從 source pptx 抽出每張投影片的文字，
    回傳 list of Dict，每個 Dict 包含：
      type     : 'title' | 'closing' | 'bullets' | 'content' | 'quote'
      title    : str
      subtitle : str | None
      bullets  : List[str] | None
      body     : str | None
      quote    : str | None
    """
    prs = Presentation(src_path)
    total = len(prs.slides)
    slides_data = []

    print(f'  共 {total} 張投影片')

    for i, slide in enumerate(prs.slides):
        # 收集所有 placeholder 的文字（依 placeholder idx 排序）
        texts = {}
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            texts[idx] = ph.text_frame.text.strip()

        title_text = texts.get(0, '') or texts.get(13, '')
        body_text  = texts.get(1, '') or texts.get(14, '')

        # 把 body_text 拆成行
        lines = [l.strip() for l in body_text.splitlines() if l.strip()]

        # 自動推斷 slide type
        if i == 0:
            slide_type = 'title'
        elif i == total - 1:
            slide_type = 'closing'
        elif body_text.startswith(('"', '\u201c', '"')):
            slide_type = 'quote'
        elif len(lines) >= 2 and all(len(l) < 100 for l in lines):
            slide_type = 'bullets'
        else:
            slide_type = 'content'

        entry = {
            'type'    : slide_type,
            'title'   : title_text,
            'subtitle': lines[0] if slide_type in ('title', 'closing') and lines else None,
            'bullets' : lines     if slide_type == 'bullets' else None,
            'body'    : body_text if slide_type == 'content'  else None,
            'quote'   : body_text.strip('""\u201c\u201d\u300c\u300d') if slide_type == 'quote' else None,
        }
        slides_data.append(entry)
        print(f'  Slide {i+1}/{total}: [{slide_type:8s}] {title_text[:40]}')

    return slides_data


# ══════════════════════════════════════════════
#  3. 自動偵測範本 layout
# ══════════════════════════════════════════════

def detect_layouts(template_path: str) -> Dict:
    """
    掃描範本的所有 slide layout，回傳：
    {
      'title'  : layout_index,
      'bullets': layout_index,
      'content': layout_index,
      'quote'  : layout_index,
      'closing': layout_index,
    }
    偵測邏輯：看 layout 名稱關鍵字 + placeholder 數量
    """
    prs = Presentation(template_path)
    mapping = {}

    print(f'\n  範本共 {len(prs.slide_layouts)} 個 layout：')
    for idx, layout in enumerate(prs.slide_layouts):
        name = layout.name.lower()
        ph_idxs = [ph.placeholder_format.idx for ph in layout.placeholders]
        print(f'    [{idx:02d}] "{layout.name}"  placeholders={ph_idxs}')

        if any(k in name for k in ('title slide', 'cover', '封面', 'opening', 'title, content')):
            mapping.setdefault('title', idx)
        if any(k in name for k in ('bullet', 'content', '內容', 'two content', 'text')):
            mapping.setdefault('bullets', idx)
            mapping.setdefault('content', idx)
        if any(k in name for k in ('quote', 'caption', '引言', 'picture with caption')):
            mapping.setdefault('quote', idx)
        if any(k in name for k in ('closing', 'end', 'thank', '結語', 'blank', 'section')):
            mapping.setdefault('closing', idx)

    # 補全未偵測到的 type
    total = len(prs.slide_layouts)
    fallbacks = {'title': 0, 'bullets': min(1, total-1), 'content': min(1, total-1),
                 'quote': min(2, total-1), 'closing': 0}
    for t, fb in fallbacks.items():
        if t not in mapping:
            mapping[t] = fb
            print(f'  ⚠  "{t}" 未偵測到關鍵字，使用 fallback layout[{fb}]')

    print(f'\n  對應結果：')
    for t, idx in mapping.items():
        print(f'    {t:10s} → layout[{idx:02d}] "{prs.slide_layouts[idx].name}"')

    return mapping


# ══════════════════════════════════════════════
#  4. 填入單張投影片
# ══════════════════════════════════════════════

def fill_placeholder(ph, paragraphs_data: list):
    """
    paragraphs_data: list of {text, level, align, font_size}
    """
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, pd in enumerate(paragraphs_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = pd.get('level', 0)
        if pd.get('align'):
            p.alignment = pd['align']

        run = p.add_run()
        run.text = pd['text']
        apply_style(run, pd['text'], font_size=pd.get('font_size'))


def render_slide(prs: Presentation, layout_idx: int, sd: Dict):
    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)
    ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    stype  = sd['type']

    # ── idx=0 : Title ──
    if 0 in ph_map and sd.get('title'):
        fill_placeholder(ph_map[0], [{'text': sd['title']}])

    # ── idx=1 : Body / Content ──
    if 1 not in ph_map:
        return  # 這個 layout 沒有 body placeholder，跳過

    if stype in ('title', 'closing') and sd.get('subtitle'):
        fill_placeholder(ph_map[1], [
            {'text': sd['subtitle'], 'align': PP_ALIGN.CENTER}
        ])

    elif stype == 'bullets' and sd.get('bullets'):
        fill_placeholder(ph_map[1], [
            {'text': b, 'level': 0, 'font_size': 16}
            for b in sd['bullets']
        ])

    elif stype == 'content' and sd.get('body'):
        fill_placeholder(ph_map[1], [
            {'text': sd['body'], 'font_size': 16}
        ])

    elif stype == 'quote' and sd.get('quote'):
        fill_placeholder(ph_map[1], [
            {'text': f'\u201c{sd["quote"]}\u201d',
             'align': PP_ALIGN.CENTER, 'font_size': 20}
        ])


# ══════════════════════════════════════════════
#  5. 主流程
# ══════════════════════════════════════════════

def main(
    source_pptx  : str = './input/source.pptx',
    template_pptx: str = './input/wisptype.pptx',
    output_pptx  : str = './output/result.pptx',
):
    print(f'━━━ Step 1  讀取來源 ━━━\n  {source_pptx}')
    slides_data = extract_slides(source_pptx)

    print(f'\n━━━ Step 2  偵測範本 ━━━\n  {template_pptx}')
    layout_map = detect_layouts(template_pptx)

    print(f'\n━━━ Step 3  填入內容 ━━━')
    prs = Presentation(template_pptx)

    # 清空範本內既有投影片
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    for i, sd in enumerate(slides_data):
        layout_idx = layout_map.get(sd['type'], 1)
        render_slide(prs, layout_idx, sd)
        print(f'  ✅ Slide {i+1:02d} [{sd["type"]:8s}] → layout[{layout_idx}]')

    Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(f'\n🎉 完成！輸出：{output_pptx}')


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='JS pptx → wisptype 範本填入工具')
    parser.add_argument('--src',      default='./input/source.pptx',   help='JS 生成的 pptx')
    parser.add_argument('--template', default='./input/wisptype.pptx', help='範本 pptx')
    parser.add_argument('--out',      default='./output/result.pptx',  help='輸出路徑')
    args = parser.parse_args()
    main(args.src, args.template, args.out)